"""
PPTX 생성 엔진

JSON spec 파일을 읽어 python-pptx로 프레젠테이션을 생성한다.
PPT_Design_Guide_2026 원칙을 자동 적용한다.

Usage:
    python generate_pptx.py <spec.json> <output.pptx> [choices.json]
"""

import json
import sys
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from PIL import Image
from copy import deepcopy
from lxml import etree
import math

# ─── 슬라이드 크기 (16:9) ─────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ─── 여백 ─────────────────────────────────────────────────────────
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
MARGIN_BOTTOM = Inches(0.5)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM

# ─── 요소 간 간격 ─────────────────────────────────────────────────
ELEMENT_GAP = Inches(0.2)  # body_elements 간 수직 여백

# ─── 폰트 ─────────────────────────────────────────────────────────
FONT_NAME = "Pretendard"
FONT_NAME_EA = "Pretendard"  # East Asian fallback
FONT_FALLBACKS = ["Malgun Gothic", "맑은 고딕", "Apple SD Gothic Neo", "sans-serif"]


def set_font_with_ea(run, font_name=FONT_NAME, ea_name=FONT_NAME_EA):
    """폰트에 East Asian typeface를 함께 설정한다 (한글 렌더링 안정화)."""
    run.font.name = font_name
    # oxml로 <a:ea> East Asian typeface 추가
    rPr = run._r.get_or_add_rPr()
    # 기존 <a:ea> 제거
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea_name)


# ─── 테이블 oxml 헬퍼 ─────────────────────────────────────────────

def _set_cell_border(cell, side, width_pt=0, color_hex=None):
    """셀의 특정 변(side)에 테두리를 설정한다. width_pt=0이면 선 제거."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    border_tag = {
        'top': 'a:lnT', 'bottom': 'a:lnB',
        'left': 'a:lnL', 'right': 'a:lnR'
    }[side]

    # 기존 제거
    for old in tcPr.findall(qn(border_tag)):
        tcPr.remove(old)

    ln = etree.SubElement(tcPr, qn(border_tag))
    if width_pt == 0:
        ln.set('w', '0')
        ln.set('cmpd', 'sng')
        no_fill = etree.SubElement(ln, qn('a:noFill'))
    else:
        ln.set('w', str(int(width_pt * 12700)))  # pt → EMU
        ln.set('cmpd', 'sng')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', color_hex.lstrip('#') if color_hex else '000000')


def _set_cell_margins(cell, top=Pt(8), bottom=Pt(8), left=Pt(12), right=Pt(12)):
    """셀 내부 여백(padding)을 설정한다."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set('marT', str(int(top)))
    tcPr.set('marB', str(int(bottom)))
    tcPr.set('marL', str(int(left)))
    tcPr.set('marR', str(int(right)))

# ─── 테마 정의 ─────────────────────────────────────────────────────
THEMES = {
    "dark": {
        "bg": "#1A1A2E",
        "text": "#FFFFFF",
        "subtitle": "#A0A0B8",
        "accent": "#E94560",
        "secondary": "#16213E",
        "card_bg": "#16213E",
        "table_header_bg": "#E94560",   # accent (빨강)
        "table_header_text": "#FFFFFF",
        "table_row_odd": "#1E2040",     # bg보다 약간 밝은 남색
        "table_row_even": "#2A2D52",    # 뚜렷한 대비를 위한 밝은 남색
        "table_text": "#E8E8F0",
        "table_border": "#3A3D60",
    },
    "light": {
        "bg": "#FFFFFF",
        "text": "#1A1A1A",
        "subtitle": "#6B7280",
        "accent": "#2563EB",
        "secondary": "#F3F4F6",
        "card_bg": "#F3F4F6",
        "table_header_bg": "#2563EB",
        "table_header_text": "#FFFFFF",
        "table_row_odd": "#FFFFFF",
        "table_row_even": "#EFF6FF",    # 연한 파랑
        "table_text": "#1E293B",
        "table_border": "#DBEAFE",
    },
    "minimal": {
        "bg": "#FAFAFA",
        "text": "#333333",
        "subtitle": "#888888",
        "accent": "#6366F1",
        "secondary": "#F0F0F0",
        "card_bg": "#FFFFFF",
        "table_header_bg": "#6366F1",
        "table_header_text": "#FFFFFF",
        "table_row_odd": "#FFFFFF",
        "table_row_even": "#F0F0F8",    # 연한 인디고
        "table_text": "#334155",
        "table_border": "#E2E8F0",
    },
    "consulting": {
        "bg": "#FFFFFF",
        "text": "#002F6C",
        "subtitle": "#4A5568",
        "accent": "#C41230",
        "secondary": "#F7F8FA",
        "card_bg": "#F7F8FA",
        "table_header_bg": "#002F6C",   # 네이비 (accent 대신 본문색)
        "table_header_text": "#FFFFFF",
        "table_row_odd": "#FFFFFF",
        "table_row_even": "#F0F4F8",    # 연한 블루그레이
        "table_text": "#1A365D",
        "table_border": "#CBD5E0",
    },
    "pitch": {
        "bg": "#0F0A1A",
        "text": "#FFFFFF",
        "subtitle": "#C4B5D0",
        "accent": "#7B2FBE",
        "secondary": "#1A1030",
        "card_bg": "#1A1030",
        "table_header_bg": "#7B2FBE",   # accent (보라)
        "table_header_text": "#FFFFFF",
        "table_row_odd": "#18122B",     # bg보다 약간 밝은 보라
        "table_row_even": "#2D1F4E",    # 뚜렷한 대비
        "table_text": "#E8E0F0",
        "table_border": "#4A3670",
    },
    "education": {
        "bg": "#F8F7FF",
        "text": "#1E1B4B",
        "subtitle": "#6B6B8D",
        "accent": "#F97316",
        "secondary": "#EDE9FE",
        "card_bg": "#FFFFFF",
        "table_header_bg": "#F97316",   # accent (오렌지)
        "table_header_text": "#FFFFFF",
        "table_row_odd": "#FFFFFF",
        "table_row_even": "#FFF7ED",    # 연한 오렌지
        "table_text": "#1E1B4B",
        "table_border": "#FED7AA",
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, hex_color):
    """슬라이드 배경색을 설정한다."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def add_text_with_markdown(paragraph, text, theme, font_size=Pt(18), is_body=True):
    """마크다운 스타일(**굵게**, *기울임*)을 해석하여 텍스트를 추가한다."""
    if not text:
        return

    # 간단한 마크다운 파싱 (정규표현식)
    # **bold**, *italic*, ~~strike~~, `code`
    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|~~.*?~~|`.*?`)', text)
    
    for part in parts:
        if not part: continue
        
        run = paragraph.add_run()
        run.font.size = font_size
        run.font.color.rgb = hex_to_rgb(theme["text"])
        set_font_with_ea(run)
        
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith('*') and part.endswith('*'):
            run.text = part[1:-1]
            run.font.italic = True
        elif part.startswith('~~') and part.endswith('~~'):
            run.text = part[2:-2]
            # python-pptx doesn't have native strikethrough in common API, skipping for simplicity or use oxml
        elif part.startswith('`') and part.endswith('`'):
            run.text = part[1:-1]
            run.font.name = "Consolas"
        else:
            run.text = part


# ─── 콘텐츠 오버플로우 감지 및 자동 분할 ──────────────────────────

def _estimate_element_height(elem, width_inches=11.733):
    """body_element의 추정 높이(인치)를 반환한다.

    렌더링 함수와 동일한 높이 계산 + ELEMENT_GAP(0.2")을 포함한다.
    이 값은 split_overflowing_slides()에서 분할 여부를 판단하는 데 사용된다.
    """
    elem_type = elem.get("type", "paragraph")
    gap = 0.2  # ELEMENT_GAP in inches
    # 18pt 한글 기준으로 인치당 약 5 글자
    chars_per_line = max(1, int(width_inches * 5))

    if elem_type == "heading":
        level = elem.get("level", 2)
        # h2: text(0.45) + bar(0.12) + gap / h3: text(0.4) + pad(0.05) + gap
        return (0.45 + 0.12 + gap) if level == 2 else (0.4 + 0.05 + gap)

    elif elem_type == "paragraph":
        text = elem.get("text", "")
        line_count = max(1, math.ceil(len(text) / chars_per_line))
        return 0.32 * line_count + 0.1 + gap

    elif elem_type == "bullet_list":
        items = elem.get("items", [])
        return 0.38 * len(items) + 0.15 + gap

    elif elem_type == "numbered_list":
        items = elem.get("items", [])
        return 0.38 * len(items) + 0.15 + gap

    elif elem_type == "blockquote":
        text = elem.get("text", "")
        line_count = max(1, math.ceil(len(text) / int(chars_per_line * 0.85)))
        return 0.35 * line_count + 0.3 + gap

    elif elem_type == "divider":
        return 0.3 + gap

    elif elem_type == "code_block":
        code = elem.get("code", "")
        line_count = len(code.split('\n'))
        return 0.25 * min(line_count, 12) + 0.4 + gap

    elif elem_type == "inline_table":
        rows = elem.get("rows", [])
        return 0.35 * (len(rows) + 1) + gap

    return 0.3 + gap  # default


def _get_layout_dimensions(layout):
    """레이아웃별 body_elements의 (가용 높이, 가용 너비) 인치를 반환한다.

    보수적으로 추정하여 실제 렌더링 시 콘텐츠가 넘치지 않도록 한다.
    타이틀 영역(~0.85") + 상하 마진을 고려한 값이다.
    """
    content_w = 11.733  # CONTENT_WIDTH in inches
    content_h = 5.0     # 보수적 가용 높이 (타이틀 + 여유 마진 고려)

    if layout == "content":
        return content_h, content_w
    elif layout == "content-image":
        return content_h, content_w * 0.45  # 좌측 45% (좁은 영역)
    elif layout == "comparison":
        return 4.2, (content_w - 0.6) / 2
    elif layout == "two-images":
        return 1.5, content_w
    else:
        return content_h, content_w


def _elements_to_plain_text(elements):
    """body_elements를 플레인 텍스트로 변환한다 (body 폴백용)."""
    lines = []
    for elem in elements:
        t = elem.get("type", "paragraph")
        if t == "heading":
            lines.append(f"▸ {elem.get('text', '')}")
        elif t == "paragraph":
            lines.append(elem.get("text", ""))
        elif t in ("bullet_list", "numbered_list"):
            prefix = "•" if t == "bullet_list" else "1."
            for i, item in enumerate(elem.get("items", [])):
                text = item if isinstance(item, str) else item.get("text", "")
                p = f"{i+1}." if t == "numbered_list" else "•"
                lines.append(f" {p} {text}")
        elif t == "blockquote":
            lines.append(f'"{elem.get("text", "")}"')
        elif t == "divider":
            lines.append("───────────────────")
        elif t == "code_block":
            lines.append(elem.get("code", ""))
        elif t == "inline_table":
            headers = elem.get("headers", [])
            if headers:
                lines.append(" | ".join(str(h) for h in headers))
            for row in elem.get("rows", []):
                lines.append(" | ".join(str(c) for c in row))
        lines.append("")
    return "\n".join(lines).strip()


def _split_by_h2_sections(body_elements):
    """body_elements를 h2 heading 기준으로 섹션 그룹으로 나눈다.

    각 섹션은 (h2_title, elements) 튜플이다.
    첫 번째 h2 이전의 요소들은 title=None으로 묶인다.
    """
    sections = []
    current_title = None
    current_elems = []

    for elem in body_elements:
        if elem.get("type") == "heading" and elem.get("level") == 2:
            if current_elems:
                sections.append((current_title, current_elems))
            current_title = elem.get("text", "")
            current_elems = []
        else:
            current_elems.append(elem)

    if current_elems:
        sections.append((current_title, current_elems))

    return sections


def _strip_md_markers(text):
    """마크다운 굵게/기울임/코드 마커를 제거한다."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'~~(.+?)~~', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    return text.strip()


def _generate_title_from_elements(elements):
    """요소들의 내용으로부터 슬라이드 제목을 자동 생성한다.

    h3 heading > paragraph 첫 구절 > 리스트 키워드 > blockquote 첫 구절
    순으로 우선 탐색하여 30자 이내의 깔끔한 제목을 만든다.
    """
    for elem in elements:
        t = elem.get("type", "")

        if t == "heading":
            return _strip_md_markers(elem.get("text", ""))[:30]

        if t == "paragraph":
            text = _strip_md_markers(elem.get("text", ""))
            # 첫 문장 추출
            for sep in ["다.", "한다.", "이다.", ". "]:
                idx = text.find(sep)
                if idx > 0 and idx < 40:
                    return text[:idx + len(sep)].strip()
            return text[:30].strip()

        if t in ("bullet_list", "numbered_list"):
            items = elem.get("items", [])
            if items:
                first = items[0] if isinstance(items[0], str) else items[0].get("text", "")
                first = _strip_md_markers(first)
                # "키워드: 설명" 패턴에서 키워드만 추출
                if ": " in first:
                    return first.split(": ")[0][:30].strip()
                if " — " in first:
                    return first.split(" — ")[0][:30].strip()
                if " → " in first:
                    return first.split(" → ")[0][:30].strip()
                return first[:30].strip()

        if t == "blockquote":
            text = _strip_md_markers(elem.get("text", ""))
            for sep in ["다.", "한다.", "이다.", ". "]:
                idx = text.find(sep)
                if idx > 0 and idx < 40:
                    return text[:idx + len(sep)].strip()
            return text[:30].strip()

        if t == "inline_table":
            headers = elem.get("headers", [])
            if headers:
                return " · ".join(str(h) for h in headers[:4])[:30]

    return ""


def _subsplit_by_height(elements, available_h, available_w):
    """요소들을 높이 기준으로 그룹으로 나눈다. 각 그룹은 available_h 이내."""
    groups = []
    current_group = []
    current_h = 0.0

    for elem in elements:
        elem_h = _estimate_element_height(elem, available_w)
        if current_h + elem_h > available_h and current_group:
            groups.append(current_group)
            current_group = []
            current_h = 0.0
        current_group.append(elem)
        current_h += elem_h

    if current_group:
        groups.append(current_group)

    return groups


def split_overflowing_slides(slides):
    """body_elements가 슬라이드 높이를 초과하는 경우 분할한다.

    분할 전략:
    1단계: h2 heading 경계에서 슬라이드를 나눈다.
    2단계: h2 섹션 하나가 여전히 초과하면, 높이 기반으로 추가 분할한다.
    3단계: 분할 슬라이드 제목 결정:
       - h2가 있으면 h2 텍스트 사용
       - 없으면 내용에서 자동 생성
    """
    result = []
    skip_layouts = {"title", "section", "closing", "image-full", "kpi", "timeline"}

    for slide_data in slides:
        layout = slide_data.get("layout", "content")
        body_elements = slide_data.get("body_elements")

        if not body_elements or layout in skip_layouts:
            result.append(slide_data)
            continue

        available_h, available_w = _get_layout_dimensions(layout)
        # content 레이아웃 기준 가용 높이 (분할 후 이미지 없는 슬라이드용)
        content_h, content_w = _get_layout_dimensions("content")

        total_h = sum(_estimate_element_height(e, available_w) for e in body_elements)

        if total_h <= available_h:
            result.append(slide_data)
            continue

        # ── 1단계: h2 기준 섹션 분리 ──
        sections = _split_by_h2_sections(body_elements)

        # ── 2단계: 섹션들을 슬라이드 그룹으로 배분 ──
        # slide_groups: [(title, [elements], is_first_slide)]
        original_title = slide_data.get("title", "")
        slide_groups = []
        current_title = original_title
        current_elems = []
        current_h = 0.0
        is_first = True

        for sec_title, sec_elems in sections:
            # 이 섹션이 들어갈 슬라이드의 가용 높이 결정
            target_h = available_h if is_first and not current_elems else content_h
            target_w = available_w if is_first and not current_elems else content_w
            sec_h = sum(_estimate_element_height(e, target_w) for e in sec_elems)

            # 현재 슬라이드에 추가하면 초과 & 이미 내용이 있으면 → 새 슬라이드
            if current_h + sec_h > target_h and current_elems:
                slide_groups.append((current_title, current_elems, is_first))
                is_first = False
                current_title = sec_title if sec_title else _generate_title_from_elements(sec_elems)
                current_elems = list(sec_elems)
                current_h = sec_h
            else:
                # h2 제목을 body에 heading으로 포함 (같은 슬라이드에 합쳐질 때)
                if sec_title and current_elems:
                    current_elems.append({"type": "heading", "level": 2, "text": sec_title})
                    current_h += _estimate_element_height({"type": "heading", "level": 2}, target_w)
                elif sec_title and not current_elems and not is_first:
                    current_title = sec_title
                elif sec_title and not current_elems and is_first:
                    current_elems.append({"type": "heading", "level": 2, "text": sec_title})
                    current_h += _estimate_element_height({"type": "heading", "level": 2}, target_w)
                current_elems.extend(sec_elems)
                current_h += sec_h

        if current_elems:
            slide_groups.append((current_title, current_elems, is_first))

        # ── 3단계: 각 그룹이 여전히 초과하면 높이 기반 추가 분할 ──
        final_groups = []
        for group_title, group_elems, first in slide_groups:
            target_h = available_h if first else content_h
            target_w = available_w if first else content_w
            group_h = sum(_estimate_element_height(e, target_w) for e in group_elems)

            if group_h <= target_h:
                final_groups.append((group_title, group_elems, first))
            else:
                # 높이 기반 추가 분할
                sub_groups = _subsplit_by_height(group_elems, target_h, target_w)
                for sgi, sub_group in enumerate(sub_groups):
                    if sgi == 0:
                        final_groups.append((group_title, sub_group, first))
                    else:
                        sub_title = _generate_title_from_elements(sub_group)
                        if not sub_title:
                            sub_title = group_title
                        final_groups.append((sub_title, sub_group, False))

        # ── 4단계: 슬라이드 생성 ──
        for gi, (group_title, group_elems, first) in enumerate(final_groups):
            new_slide = deepcopy(slide_data)
            new_slide["title"] = group_title
            new_slide["body_elements"] = group_elems
            new_slide["body"] = _elements_to_plain_text(group_elems)

            if not first:
                if "image" in new_slide:
                    del new_slide["image"]
                if layout == "content-image":
                    new_slide["layout"] = "content"

            result.append(new_slide)

    return result


# ─── 구조화된 본문 렌더러 (body_elements) ─────────────────────────

def render_body_elements(slide, elements, theme, left, top, width, height):
    """body_elements 배열을 슬라이드에 렌더링한다.

    각 요소의 type에 따라 적절한 디자인을 적용한다:
    - heading (h2/h3): 서브 제목
    - paragraph: 일반 텍스트
    - bullet_list: 불릿 목록 (다단계)
    - numbered_list: 번호 목록
    - blockquote: 인용구 (좌측 accent 바 + 배경)
    - divider: 수평 구분선
    - code_block: 인라인 코드 블록
    - inline_table: 본문 내 소형 테이블
    """
    cursor_top = top  # 현재 y 위치 추적
    remaining = height

    for elem in elements:
        elem_type = elem.get("type", "paragraph")
        if remaining <= 0:
            break

        if elem_type == "heading":
            cursor_top, remaining = _render_heading(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "paragraph":
            cursor_top, remaining = _render_paragraph(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "bullet_list":
            cursor_top, remaining = _render_bullet_list(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "numbered_list":
            cursor_top, remaining = _render_numbered_list(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "blockquote":
            cursor_top, remaining = _render_blockquote(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "divider":
            cursor_top, remaining = _render_divider(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "code_block":
            cursor_top, remaining = _render_code_block(
                slide, elem, theme, left, cursor_top, width, remaining
            )
        elif elem_type == "inline_table":
            cursor_top, remaining = _render_inline_table(
                slide, elem, theme, left, cursor_top, width, remaining
            )

    return cursor_top


def _render_heading(slide, elem, theme, left, top, width, remaining):
    """## (h2) 또는 ### (h3) 서브 제목을 렌더링한다."""
    level = elem.get("level", 2)
    text = elem.get("text", "")

    if level == 2:
        font_size = Pt(24)
        text_h = Inches(0.45)
        box = slide.shapes.add_textbox(left, top, width, text_h)
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(theme["text"])
        set_font_with_ea(run)

        # accent 하단 라인 — 텍스트 아래에 여유 있게 배치
        bar_top = top + text_h + Inches(0.02)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, bar_top, Inches(1.5), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
        line.line.fill.background()

        box_h = text_h + Inches(0.12)  # 텍스트 + bar + 여유
    else:
        font_size = Pt(20)
        text_h = Inches(0.4)
        box = slide.shapes.add_textbox(left, top, width, text_h)
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(theme["text"])
        set_font_with_ea(run)

        box_h = text_h + Inches(0.05)

    total = box_h + ELEMENT_GAP
    return top + total, remaining - total


def _render_paragraph(slide, elem, theme, left, top, width, remaining):
    """일반 텍스트 단락을 렌더링한다."""
    text = elem.get("text", "")
    if not text.strip():
        gap = Inches(0.15)
        return top + gap, remaining - gap

    # 너비 기준 줄 수 추정 (인치당 ~5 한글 글자)
    chars_per_line = max(1, int(width / Inches(1) * 5))
    line_count = max(1, math.ceil(len(text) / chars_per_line))
    box_h = Inches(0.32 * line_count + 0.1)

    box = slide.shapes.add_textbox(left, top, width, box_h)
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    add_text_with_markdown(para, text, theme, font_size=Pt(18))

    total = box_h + ELEMENT_GAP
    return top + total, remaining - total


def _render_bullet_list(slide, elem, theme, left, top, width, remaining):
    """불릿 목록을 렌더링한다 (다단계 지원)."""
    items = elem.get("items", [])
    base_level = elem.get("level", 0)
    if not items:
        return top, remaining

    indent_per_level = Inches(0.3)
    line_h = Inches(0.38)  # 항목 간 충분한 간격
    total_h = line_h * len(items) + Inches(0.15)

    box = slide.shapes.add_textbox(left, top, width, total_h)
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        # 아이템이 dict이면 중첩 레벨 지원
        if isinstance(item, dict):
            item_text = item.get("text", "")
            item_level = item.get("level", base_level)
        else:
            item_text = str(item)
            item_level = base_level

        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = item_level
        para.space_before = Pt(6)
        para.space_after = Pt(4)

        # 불릿 마커
        bullet_run = para.add_run()
        bullet_run.text = "•  "
        bullet_run.font.size = Pt(18)
        bullet_run.font.color.rgb = hex_to_rgb(theme["accent"])
        set_font_with_ea(bullet_run)

        add_text_with_markdown(para, item_text, theme, font_size=Pt(18))

    total = total_h + ELEMENT_GAP
    return top + total, remaining - total


def _render_numbered_list(slide, elem, theme, left, top, width, remaining):
    """번호 목록을 렌더링한다."""
    items = elem.get("items", [])
    if not items:
        return top, remaining

    line_h = Inches(0.38)
    total_h = line_h * len(items) + Inches(0.15)

    box = slide.shapes.add_textbox(left, top, width, total_h)
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(6)
        para.space_after = Pt(4)

        # 번호
        num_run = para.add_run()
        num_run.text = f"{i + 1}.  "
        num_run.font.size = Pt(18)
        num_run.font.bold = True
        num_run.font.color.rgb = hex_to_rgb(theme["accent"])
        set_font_with_ea(num_run)

        item_text = item.get("text", item) if isinstance(item, dict) else str(item)
        add_text_with_markdown(para, item_text, theme, font_size=Pt(18))

    total = total_h + ELEMENT_GAP
    return top + total, remaining - total


def _render_blockquote(slide, elem, theme, left, top, width, remaining):
    """인용구를 렌더링한다 (좌측 accent 바 + secondary 배경)."""
    text = elem.get("text", "")
    if not text.strip():
        return top, remaining

    line_count = max(1, len(text) // 55 + 1)
    box_h = Inches(0.35 * line_count + 0.3)
    padding = Inches(0.15)

    # 배경 박스
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, box_h
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(theme["secondary"])
    bg.line.fill.background()

    # 좌측 accent 바
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, Pt(5), box_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    # 텍스트
    txt = slide.shapes.add_textbox(
        left + Inches(0.3), top + padding, width - Inches(0.45), box_h - padding * 2
    )
    tf = txt.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    add_text_with_markdown(para, text, theme, font_size=Pt(16))
    # 이탤릭 스타일 적용
    for run in para.runs:
        run.font.italic = True

    total = box_h + ELEMENT_GAP
    return top + total, remaining - total


def _render_divider(slide, elem, theme, left, top, width, remaining):
    """수평 구분선을 렌더링한다."""
    gap = Inches(0.15)
    line_top = top + gap

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.5), line_top, width - Inches(1.0), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme["subtitle"])
    line.line.fill.background()

    total_h = gap * 2 + Pt(1.5) + ELEMENT_GAP
    return top + total_h, remaining - total_h


def _render_code_block(slide, elem, theme, left, top, width, remaining):
    """인라인 코드 블록을 렌더링한다 (어두운 배경 + 모노스페이스)."""
    code = elem.get("code", "")
    language = elem.get("language", "")
    if not code.strip():
        return top, remaining

    lines = code.split('\n')
    line_count = len(lines)
    box_h = Inches(0.25 * min(line_count, 12) + 0.4)
    padding = Inches(0.12)

    # 어두운 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, box_h
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
    bg.line.color.rgb = hex_to_rgb(theme["accent"])
    bg.line.width = Pt(0.5)

    # 언어 라벨
    if language:
        label_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Pt(4), Inches(1.5), Inches(0.25)
        )
        lf = label_box.text_frame
        lp = lf.paragraphs[0]
        lr = lp.add_run()
        lr.text = language
        lr.font.size = Pt(10)
        lr.font.color.rgb = hex_to_rgb(theme["accent"])
        lr.font.name = "Consolas"
        code_top_offset = Inches(0.3)
    else:
        code_top_offset = padding

    # 코드 텍스트
    code_box = slide.shapes.add_textbox(
        left + Inches(0.15), top + code_top_offset,
        width - Inches(0.3), box_h - code_top_offset - padding
    )
    tf = code_box.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = code
    run.font.name = "Consolas"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(220, 220, 220)

    total = box_h + ELEMENT_GAP
    return top + total, remaining - total


def _render_inline_table(slide, elem, theme, left, top, width, remaining):
    """본문 내 소형 테이블을 렌더링한다."""
    headers = elem.get("headers", [])
    rows = elem.get("rows", [])
    if not headers:
        return top, remaining

    rows_count = len(rows) + 1
    cols_count = len(headers)
    row_h = Inches(0.35)
    table_h = row_h * rows_count

    table_shape = slide.shapes.add_table(
        rows_count, cols_count,
        left, top, width, table_h
    )
    table = table_shape.table

    # 테이블 전용 색상 (테마에 정의되어 있으면 사용, 없으면 폴백)
    t_header_bg = theme.get("table_header_bg", theme["accent"])
    t_header_text = theme.get("table_header_text", "#FFFFFF")
    t_row_odd = theme.get("table_row_odd", theme["bg"])
    t_row_even = theme.get("table_row_even", theme["secondary"])
    t_text = theme.get("table_text", theme["text"])
    t_border = theme.get("table_border", theme.get("subtitle", "#CCCCCC"))

    # 헤더
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(t_header_bg)
        _set_cell_border(cell, 'left', 0)
        _set_cell_border(cell, 'right', 0)
        _set_cell_border(cell, 'top', 0)
        _set_cell_border(cell, 'bottom', width_pt=2, color_hex=t_border)
        _set_cell_margins(cell, top=Pt(6), bottom=Pt(6), left=Pt(10), right=Pt(10))
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(h)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(t_header_text)
        set_font_with_ea(run)

    # 데이터 행
    for r, row_data in enumerate(rows):
        row_bg = t_row_even if r % 2 == 0 else t_row_odd
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(row_bg)
            _set_cell_border(cell, 'left', 0)
            _set_cell_border(cell, 'right', 0)
            _set_cell_border(cell, 'top', 0)
            _set_cell_border(cell, 'bottom', width_pt=0.5, color_hex=t_border)
            _set_cell_margins(cell, top=Pt(5), bottom=Pt(5), left=Pt(10), right=Pt(10))
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(val)
            run.font.size = Pt(13)
            run.font.color.rgb = hex_to_rgb(t_text)
            set_font_with_ea(run)

    total_h = table_h + ELEMENT_GAP
    return top + total_h, remaining - total_h


# ─── 슬라이드 제목 렌더 헬퍼 ──────────────────────────────────────

def _render_slide_title(slide, title_text, theme, width=None):
    """슬라이드 제목을 렌더링하고 body 시작 y 위치를 반환한다.

    제목이 길어 2줄 이상이 되면 title 영역을 동적으로 확장하여
    accent bar나 본문 콘텐츠와 겹치지 않도록 한다.
    """
    if width is None:
        width = CONTENT_WIDTH

    # 32pt 폰트 기준 제목 줄 수 추정 (인치당 ~4.5 한글 글자)
    chars_per_line = max(1, int(width / Inches(1) * 4.5))
    line_count = max(1, math.ceil(len(title_text) / chars_per_line))
    title_h = Inches(0.5 * line_count + 0.1)

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, width, title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = title_text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    body_top = MARGIN_TOP + title_h + Inches(0.25)
    return body_top


# ─── 레이아웃 핸들러 ──────────────────────────────────────────────

def layout_title(prs, slide_data, theme):
    """제목 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, theme["bg"])

    title_text = slide_data.get("title", "")

    # 제목 줄 수 추정 (54pt Bold 한글 기준)
    # 54pt 한글 글자 ≈ 0.75" 너비 → 인치당 약 1.3자
    chars_per_line = max(1, int(CONTENT_WIDTH / Inches(1) * 1.3))
    title_lines = max(1, math.ceil(len(title_text) / chars_per_line))
    title_h = Inches(0.85 * title_lines)

    # 제목 시작 y — 줄 수에 따라 상단으로 올림
    if title_lines == 1:
        title_top = Inches(2.8)
    elif title_lines == 2:
        title_top = Inches(2.2)
    else:
        title_top = Inches(1.6)

    # 메인 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, title_top, CONTENT_WIDTH, title_h
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = title_text
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    # accent bar와 부제목 위치 — 제목 하단 기준으로 동적 배치
    bar_y = title_top + title_h + Inches(0.25)
    subtitle_y = bar_y + Inches(0.2)

    # 부제목
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            MARGIN_LEFT, subtitle_y, CONTENT_WIDTH, Inches(0.8)
        )
        tf2 = sub_box.text_frame
        para2 = tf2.paragraphs[0]
        para2.alignment = PP_ALIGN.CENTER
        run2 = para2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(24)
        run2.font.color.rgb = hex_to_rgb(theme["subtitle"])
        set_font_with_ea(run2)

    # 강조 데코레이션 — 제목 아래에 동적 배치
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), bar_y, Inches(2.3), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    line.line.fill.background()

    return slide


def layout_section(prs, slide_data, theme):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["secondary"])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(3.0), CONTENT_WIDTH, Inches(1.5)
    )
    tf = title_box.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    # 사이드 바 데코
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Pt(0), Inches(2.8), Inches(0.4), Inches(1.8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    return slide


def layout_content(prs, slide_data, theme):
    """기본 콘텐츠 슬라이드 (body_elements 구조화 렌더링 또는 기존 body 텍스트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 제목 (동적 높이)
    body_top = _render_slide_title(slide, slide_data.get("title", ""), theme)
    body_height = SLIDE_HEIGHT - MARGIN_BOTTOM - body_top

    # body_elements가 있으면 구조화된 렌더링, 없으면 기존 body 텍스트 방식
    body_elements = slide_data.get("body_elements")
    if body_elements:
        render_body_elements(
            slide, body_elements, theme,
            MARGIN_LEFT, body_top, CONTENT_WIDTH, body_height
        )
    else:
        body = slide_data.get("body", "")
        if body:
            content_box = slide.shapes.add_textbox(
                MARGIN_LEFT, body_top, CONTENT_WIDTH, body_height
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            lines = body.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # 불릿 포인트 처리
                clean_line = line.strip()
                if clean_line.startswith('- ') or clean_line.startswith('* '):
                    p.level = 0
                    add_text_with_markdown(p, clean_line[2:], theme)
                elif clean_line.startswith('  - ') or clean_line.startswith('  * '):
                    p.level = 1
                    add_text_with_markdown(p, clean_line[4:], theme)
                else:
                    add_text_with_markdown(p, line, theme)

    return slide


def layout_content_image(prs, slide_data, theme):
    """좌측 텍스트, 우측 이미지 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 제목 (동적 높이)
    body_top = _render_slide_title(slide, slide_data.get("title", ""), theme)

    # 텍스트 영역 (좌측 45%)
    body_width = CONTENT_WIDTH * 0.45
    body_height = SLIDE_HEIGHT - MARGIN_BOTTOM - body_top

    body_elements = slide_data.get("body_elements")
    if body_elements:
        render_body_elements(
            slide, body_elements, theme,
            MARGIN_LEFT, body_top, body_width, body_height
        )
    else:
        content_box = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, body_width, body_height
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        body = slide_data.get("body", "")
        for i, line in enumerate(body.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            add_text_with_markdown(p, line, theme, font_size=Pt(18))

    # 이미지 영역 (우측 50%)
    img_path = slide_data.get("image")
    if img_path and os.path.exists(img_path):
        img_left = MARGIN_LEFT + body_width + Inches(0.4)
        img_top = body_top
        img_width = CONTENT_WIDTH * 0.5
        img_height = SLIDE_HEIGHT - MARGIN_BOTTOM - body_top
        
        # 가로세로 비율 유지하며 채우기
        slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)

    return slide


def layout_image_full(prs, slide_data, theme):
    """전체 화면 이미지 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "#000000")

    img_path = slide_data.get("image")
    if img_path and os.path.exists(img_path):
        # 꽉 채우기
        slide.shapes.add_picture(img_path, 0, 0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT)

    # 하단 캡션/제목
    title = slide_data.get("title", "")
    if title:
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(1.2), SLIDE_WIDTH, Inches(1.2)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.4
        overlay.line.fill.background()

        txt = slide.shapes.add_textbox(Inches(0.5), SLIDE_HEIGHT - Inches(1.0), SLIDE_WIDTH - Inches(1.0), Inches(0.8))
        para = txt.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        set_font_with_ea(run)

    return slide


def layout_two_images(prs, slide_data, theme):
    """텍스트 + 2개 이미지 병렬 배치"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 제목
    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    # 텍스트 (상단 30%)
    body_elements = slide_data.get("body_elements")
    if body_elements:
        render_body_elements(
            slide, body_elements, theme,
            MARGIN_LEFT, MARGIN_TOP + Inches(0.9), CONTENT_WIDTH, Inches(1.5)
        )
    else:
        body = slide_data.get("body", "")
        if body:
            txt_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP + Inches(0.9), CONTENT_WIDTH, Inches(1.5))
            tf = txt_box.text_frame; tf.word_wrap = True
            for i, line in enumerate(body.split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                add_text_with_markdown(p, line, theme, font_size=Pt(16))

    # 이미지 2개 (하단 60%)
    imgs = slide_data.get("images", [])
    if len(imgs) >= 2:
        w = (CONTENT_WIDTH - Inches(0.4)) / 2
        h = Inches(3.5)
        top = SLIDE_HEIGHT - h - MARGIN_BOTTOM
        if os.path.exists(imgs[0]):
            slide.shapes.add_picture(imgs[0], MARGIN_LEFT, top, width=w)
        if os.path.exists(imgs[1]):
            slide.shapes.add_picture(imgs[1], MARGIN_LEFT + w + Inches(0.4), top, width=w)

    return slide


def layout_grid_images(prs, slide_data, theme):
    """이미지 그리드 레이아웃 (2x2 등)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    imgs = slide_data.get("images", [])
    if not imgs: return slide

    n = len(imgs)
    cols = 2 if n > 1 else 1
    rows = (n + 1) // 2
    
    gap = Inches(0.2)
    w = (CONTENT_WIDTH - (cols-1)*gap) / cols
    h = (CONTENT_HEIGHT - Inches(1.0) - (rows-1)*gap) / rows
    
    for i, img_path in enumerate(imgs):
        if not os.path.exists(img_path): continue
        r, c = i // cols, i % cols
        slide.shapes.add_picture(
            img_path, 
            MARGIN_LEFT + c*(w+gap), 
            MARGIN_TOP + Inches(1.0) + r*(h+gap),
            width=w
        )
    return slide

def layout_text_left_img_right(prs, slide_data, theme):
    """텍스트(좌) + 이미지 그리드(우)"""
    return layout_content_image(prs, slide_data, theme) # 공용으로 사용하거나 커스텀 구현

def layout_img_left_text_right(prs, slide_data, theme):
    """이미지 그리드(좌) + 텍스트(우)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])
    
    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", ""); run.font.size = Pt(32); run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["text"]); set_font_with_ea(run)

    img_width = CONTENT_WIDTH * 0.5
    imgs = slide_data.get("images", [])
    if imgs and os.path.exists(imgs[0]):
        slide.shapes.add_picture(imgs[0], MARGIN_LEFT, MARGIN_TOP + Inches(1.0), width=img_width)

    body_left = MARGIN_LEFT + img_width + Inches(0.4)
    text_width = CONTENT_WIDTH - img_width - Inches(0.4)
    text_top = MARGIN_TOP + Inches(1.0)
    text_height = CONTENT_HEIGHT - Inches(1.0)

    body_elements = slide_data.get("body_elements")
    if body_elements:
        render_body_elements(slide, body_elements, theme, body_left, text_top, text_width, text_height)
    else:
        content_box = slide.shapes.add_textbox(body_left, text_top, text_width, text_height)
        tf = content_box.text_frame; tf.word_wrap = True
        body = slide_data.get("body", "")
        for i, line in enumerate(body.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            add_text_with_markdown(p, line, theme, font_size=Pt(18))
    return slide

def layout_text_top_img_bottom(prs, slide_data, theme):
    """텍스트(상) + 이미지 가로 나열(하)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])
    # 제목
    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", ""); run.font.size = Pt(28); run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["text"]); set_font_with_ea(run)
    # 본문
    txt_h = Inches(1.5)
    body_elements = slide_data.get("body_elements")
    if body_elements:
        render_body_elements(
            slide, body_elements, theme,
            MARGIN_LEFT, MARGIN_TOP + Inches(0.9), CONTENT_WIDTH, txt_h
        )
    else:
        body = slide_data.get("body", "")
        txt_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP + Inches(0.9), CONTENT_WIDTH, txt_h)
        tf = txt_box.text_frame; tf.word_wrap = True
        for i, line in enumerate(body.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            add_text_with_markdown(p, line, theme, font_size=Pt(16))
    # 이미지들
    imgs = slide_data.get("images", [])
    if imgs:
        n = len(imgs)
        gap = Inches(0.2)
        w = (CONTENT_WIDTH - (n-1)*gap) / n
        top = MARGIN_TOP + Inches(0.9) + txt_h + Inches(0.3)
        for i, img in enumerate(imgs):
            if os.path.exists(img):
                slide.shapes.add_picture(img, MARGIN_LEFT + i*(w+gap), top, width=w)
    return slide


def layout_comparison(prs, slide_data, theme):
    """좌우 비교 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    col_w = (CONTENT_WIDTH - Inches(0.6)) / 2
    
    def render_col(data, x):
        col_top = MARGIN_TOP + Inches(1.2)
        col_height = CONTENT_HEIGHT - Inches(1.2)

        # 컬럼 제목
        title_box = slide.shapes.add_textbox(x, col_top, col_w, Inches(0.4))
        p0 = title_box.text_frame.paragraphs[0]
        r0 = p0.add_run()
        r0.text = data.get("title", "")
        r0.font.size = Pt(22); r0.font.bold = True; r0.font.color.rgb = hex_to_rgb(theme["accent"])
        set_font_with_ea(r0)

        # 컬럼 내용
        body_top = col_top + Inches(0.5)
        body_h = col_height - Inches(0.5)
        col_elements = data.get("body_elements")
        if col_elements:
            render_body_elements(slide, col_elements, theme, x, body_top, col_w, body_h)
        else:
            box = slide.shapes.add_textbox(x, body_top, col_w, body_h)
            tf = box.text_frame; tf.word_wrap = True
            for i, line in enumerate(data.get("body", "").split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                add_text_with_markdown(p, line, theme, font_size=Pt(16))

    render_col(slide_data.get("left", {}), MARGIN_LEFT)
    render_col(slide_data.get("right", {}), MARGIN_LEFT + col_w + Inches(0.6))

    return slide


def layout_table(prs, slide_data, theme):
    """데이터 테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    if not headers: return slide

    rows_count = len(rows) + 1
    cols_count = len(headers)
    
    table_shape = slide.shapes.add_table(
        rows_count, cols_count, 
        MARGIN_LEFT, MARGIN_TOP + Inches(1.2), 
        CONTENT_WIDTH, Inches(0.5) * rows_count
    )
    table = table_shape.table

    # 테이블 전용 색상
    t_header_bg = theme.get("table_header_bg", theme["accent"])
    t_header_text = theme.get("table_header_text", "#FFFFFF")
    t_row_odd = theme.get("table_row_odd", theme["bg"])
    t_row_even = theme.get("table_row_even", theme["secondary"])
    t_text = theme.get("table_text", theme["text"])
    t_border = theme.get("table_border", theme.get("subtitle", "#CCCCCC"))

    # 헤더 스타일
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(t_header_bg)
        _set_cell_border(cell, 'bottom', width_pt=2, color_hex=t_border)
        _set_cell_margins(cell, top=Pt(6), bottom=Pt(6), left=Pt(10), right=Pt(10))
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(h)
        run.font.size = Pt(18); run.font.bold = True
        run.font.color.rgb = hex_to_rgb(t_header_text)
        set_font_with_ea(run)

    # 데이터 행
    for r, row_data in enumerate(rows):
        row_bg = t_row_even if r % 2 == 0 else t_row_odd
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(row_bg)
            _set_cell_border(cell, 'left', 0)
            _set_cell_border(cell, 'right', 0)
            _set_cell_border(cell, 'top', 0)
            _set_cell_border(cell, 'bottom', width_pt=0.5, color_hex=t_border)
            _set_cell_margins(cell, top=Pt(5), bottom=Pt(5), left=Pt(10), right=Pt(10))
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(val)
            run.font.size = Pt(16)
            run.font.color.rgb = hex_to_rgb(t_text)
            set_font_with_ea(run)

    return slide


def layout_code(prs, slide_data, theme):
    """코드 블록 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    # 코드 배경 박스
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, MARGIN_TOP + Inches(1.0),
        CONTENT_WIDTH, CONTENT_HEIGHT - Inches(1.0)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30, 30, 30) # 항상 어두운 배경
    bg.line.color.rgb = hex_to_rgb(theme["accent"])

    code_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(1.1),
        CONTENT_WIDTH - Inches(0.4), CONTENT_HEIGHT - Inches(1.2)
    )
    tf = code_box.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = slide_data.get("code", "")
    run.font.name = "Consolas"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(220, 220, 220)

    return slide


def layout_kpi(prs, slide_data, theme):
    """KPI / 지표 레이아웃 (카드형)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    title_box.text_frame.paragraphs[0].text = slide_data.get("title", "")

    metrics = slide_data.get("metrics", [])
    if not metrics: return slide

    card_count = len(metrics)
    card_w = (CONTENT_WIDTH - (card_count-1)*Inches(0.3)) / card_count
    card_h = Inches(2.5)
    
    for i, m in enumerate(metrics):
        left = MARGIN_LEFT + i * (card_w + Inches(0.3))
        top = Inches(3.0)
        
        # 카드 배경
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        rect.fill.solid(); rect.fill.fore_color.rgb = hex_to_rgb(theme["card_bg"])
        rect.line.color.rgb = hex_to_rgb(theme["secondary"])
        
        # 값
        v_box = slide.shapes.add_textbox(left, top + Inches(0.5), card_w, Inches(0.8))
        p_v = v_box.text_frame.paragraphs[0]; p_v.alignment = PP_ALIGN.CENTER
        r_v = p_v.add_run(); r_v.text = str(m.get("value", ""))
        r_v.font.size = Pt(44); r_v.font.bold = True; r_v.font.color.rgb = hex_to_rgb(theme["accent"])
        
        # 라벨
        l_box = slide.shapes.add_textbox(left, top + Inches(1.4), card_w, Inches(0.5))
        p_l = l_box.text_frame.paragraphs[0]; p_l.alignment = PP_ALIGN.CENTER
        r_l = p_l.add_run(); r_l.text = str(m.get("label", ""))
        r_l.font.size = Pt(18); r_l.font.color.rgb = hex_to_rgb(theme["text"])

    return slide


def layout_timeline(prs, slide_data, theme):
    """타임라인 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_box = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = slide_data.get("title", "")
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    events = slide_data.get("events", [])
    if not events: return slide

    # 가로 중앙 선
    line_y = SLIDE_HEIGHT / 2 + Inches(0.5)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_LEFT, line_y, CONTENT_WIDTH, Pt(2)
    )
    line.fill.solid(); line.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    line.line.fill.background()

    count = len(events)
    step = CONTENT_WIDTH / count
    
    for i, ev in enumerate(events):
        x = MARGIN_LEFT + i * step + (step / 2)
        
        # 원 점
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Pt(6), line_y - Pt(6), Pt(12), Pt(12))
        dot.fill.solid(); dot.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
        dot.line.fill.background()
        
        # 날짜 (선 위)
        d_box = slide.shapes.add_textbox(x - Inches(0.8), line_y - Inches(0.6), Inches(1.6), Inches(0.4))
        p_d = d_box.text_frame.paragraphs[0]; p_d.alignment = PP_ALIGN.CENTER
        r_d = p_d.add_run(); r_d.text = ev.get("date", "")
        r_d.font.size = Pt(18); r_d.font.bold = True; r_d.font.color.rgb = hex_to_rgb(theme["accent"])
        
        # 설명 (선 아래)
        desc_box = slide.shapes.add_textbox(x - Inches(1.0), line_y + Inches(0.3), Inches(2.0), Inches(1.0))
        tf = desc_box.text_frame; tf.word_wrap = True
        p_desc = tf.paragraphs[0]; p_desc.alignment = PP_ALIGN.CENTER
        r_desc = p_desc.add_run(); r_desc.text = ev.get("description", ev.get("title", ""))
        r_desc.font.size = Pt(14); r_desc.font.color.rgb = hex_to_rgb(theme["text"])

    return slide


def layout_closing(prs, slide_data, theme):
    """클로징 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    title_text = slide_data.get("title", "감사합니다")

    # 제목 줄 수 추정 (54pt Bold 한글 기준, 인치당 ~1.3자)
    chars_per_line = max(1, int(CONTENT_WIDTH / Inches(1) * 1.3))
    title_lines = max(1, math.ceil(len(title_text) / chars_per_line))
    title_h = Inches(0.85 * title_lines)
    if title_lines == 1:
        title_top = Inches(2.8)
    elif title_lines == 2:
        title_top = Inches(2.2)
    else:
        title_top = Inches(1.6)

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, title_top, CONTENT_WIDTH, title_h
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = title_text
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(theme["text"])
    set_font_with_ea(run)

    # bar/부제목/연락처 — 제목 하단 기준 동적 배치
    bar_y = title_top + title_h + Inches(0.25)
    subtitle_y = bar_y + Inches(0.2)

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            MARGIN_LEFT, subtitle_y, CONTENT_WIDTH, Inches(0.8)
        )
        tf2 = sub_box.text_frame
        para2 = tf2.paragraphs[0]
        para2.alignment = PP_ALIGN.CENTER
        run2 = para2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(22)
        run2.font.color.rgb = hex_to_rgb(theme["subtitle"])
        set_font_with_ea(run2)
        contact_y = subtitle_y + Inches(0.9)
    else:
        contact_y = subtitle_y + Inches(0.3)

    contact = slide_data.get("contact", "")
    if contact:
        c_box = slide.shapes.add_textbox(
            MARGIN_LEFT, contact_y, CONTENT_WIDTH, Inches(0.5)
        )
        tf3 = c_box.text_frame
        para3 = tf3.paragraphs[0]
        para3.alignment = PP_ALIGN.CENTER
        run3 = para3.add_run()
        run3.text = contact
        run3.font.size = Pt(16)
        run3.font.color.rgb = hex_to_rgb(theme["subtitle"])
        set_font_with_ea(run3)

    # 강조 라인 — 제목 하단 기준 동적 배치
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), bar_y, Inches(2.3), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    line.line.fill.background()

    return slide


# ─── 레이아웃 디스패처 ──────────────────────────────────────────────

LAYOUT_HANDLERS = {
    "title": layout_title,
    "section": layout_section,
    "content": layout_content,
    "content-image": layout_content_image,
    "image-full": layout_image_full,
    "two-images": layout_two_images,
    "grid-images": layout_grid_images,
    "text-left-img-right": layout_text_left_img_right,
    "img-left-text-right": layout_img_left_text_right,
    "text-top-img-bottom": layout_text_top_img_bottom,
    "comparison": layout_comparison,
    "table": layout_table,
    "code": layout_code,
    "kpi": layout_kpi,
    "timeline": layout_timeline,
    "closing": layout_closing,
}


def add_speaker_notes(slide, notes: str):
    """슬라이드에 발표자 노트를 추가한다."""
    if notes:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes


def resolve_image_paths(spec, base_dir):
    """슬라이드의 이미지 경로를 spec.json 기준 절대경로로 변환한다."""
    base_dir = os.path.abspath(base_dir)
    for slide_data in spec.get("slides", []):
        # 단일 이미지
        if "image" in slide_data and slide_data["image"]:
            img = slide_data["image"]
            if not os.path.isabs(img):
                slide_data["image"] = os.path.abspath(os.path.join(base_dir, img))
            else:
                slide_data["image"] = os.path.abspath(img)
        # 다중 이미지
        if "images" in slide_data and slide_data["images"]:
            resolved = []
            for img in slide_data["images"]:
                if not os.path.isabs(img):
                    resolved.append(os.path.abspath(os.path.join(base_dir, img)))
                else:
                    resolved.append(os.path.abspath(img))
            slide_data["images"] = resolved


def generate(spec_path: str, output_path: str, choices_path: str = None):
    """JSON spec으로부터 PPTX를 생성한다."""
    with open(spec_path, "r", encoding="utf-8") as f:
        spec = json.load(f)

    choices = {}
    if choices_path and os.path.exists(choices_path):
        with open(choices_path, "r", encoding="utf-8") as f:
            choices = json.load(f)
        print(f"choices.json 로드 완료: {choices_path}")

    # 이미지 경로를 spec.json 기준으로 해석
    spec_dir = os.path.dirname(os.path.abspath(spec_path))
    resolve_image_paths(spec, spec_dir)

    meta = spec.get("meta", {})
    # choices.json의 테마가 우선
    theme_name = choices.get("theme", meta.get("theme", "light"))
    theme = THEMES.get(theme_name, THEMES["light"])

    # 콘텐츠 오버플로우 슬라이드 자동 분할
    original_count = len(spec.get("slides", []))
    spec["slides"] = split_overflowing_slides(spec.get("slides", []))
    split_count = len(spec["slides"])
    if split_count > original_count:
        print(f"  콘텐츠 오버플로우 감지: {original_count}장 → {split_count}장으로 자동 분할")

    # 슬라이드별 레이아웃 오버라이드 적용
    slide_overrides = choices.get("slide_overrides", {})

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, slide_data in enumerate(spec.get("slides", [])):
        # choices.json에 해당 슬라이드 오버라이드가 있으면 반영
        if str(i) in slide_overrides:
            override = slide_overrides[str(i)]
            if "layout" in override:
                slide_data["layout"] = override["layout"]

        layout = slide_data.get("layout", "content")
        handler = LAYOUT_HANDLERS.get(layout, layout_content)
        slide = handler(prs, slide_data, theme)

        notes = slide_data.get("notes", "")
        if notes:
            add_speaker_notes(slide, notes)

    prs.save(output_path)
    print(f"PPTX 생성 완료: {output_path}")
    print(f"  테마: {theme_name}")
    print(f"  슬라이드 수: {len(spec.get('slides', []))}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generate_pptx.py <spec.json> <output.pptx> [choices.json]")
        sys.exit(1)

    spec_arg = sys.argv[1]
    output_arg = sys.argv[2]
    choices_arg = sys.argv[3] if len(sys.argv) > 3 else None

    generate(spec_arg, output_arg, choices_arg)
